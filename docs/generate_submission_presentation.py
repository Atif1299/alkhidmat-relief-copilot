"""
Alkhidmat Relief Copilot — Hackathon Submission Deck (15 slides)
Dark theme · Bold type · High contrast · Hackathon-winning style
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Brand palette (dark hackathon) ──────────────────────────────────────
BG_DARK   = "#0a0f0d"
BG_CARD   = "#131a16"
CANOPY    = "#14b87a"   # vibrant green accent
BRASS     = "#d4a84b"   # warm accent (HITL / warnings)
SAFFRON   = "#e8732a"   # hot accent (critical / CTA)
WHITE     = "#f0f4f1"
MUTED     = "#8a9b90"
DARK_LINE = "#1e2b24"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def solid_bg(slide, hex_color: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)


def accent_bar(slide, y_in: float = 0, h_in: float = 0.06):
    """Thin accent line across top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y_in), SLIDE_W, Inches(h_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(CANOPY)
    bar.line.fill.background()


def add_text(slide, text: str, x: float, y: float, w: float, h: float,
             size: int = 18, bold: bool = False, color: str = WHITE,
             align=PP_ALIGN.LEFT, font_name: str = "Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = font_name
    p.alignment = align
    return tf


def add_multiline(slide, lines: list[str], x: float, y: float, w: float, h: float,
                  size: int = 18, color: str = WHITE, bold: bool = False,
                  bullet_color: str = CANOPY, spacing: int = 1600):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet dot
        run_dot = p.add_run()
        run_dot.text = "●  "
        run_dot.font.size = Pt(size - 4)
        run_dot.font.color.rgb = rgb(bullet_color)
        run_dot.font.name = "Segoe UI"
        # Text
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        run.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.LEFT
    return tf


def add_chip(slide, x: float, y: float, w: float, h: float,
             label: str, fill: str = CANOPY, text_color: str = BG_DARK,
             size: int = 16, bold: bool = True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    # Center text
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(text_color)
    run.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    return shape


def footer(slide):
    add_text(slide, "github.com/Atif1299/alkhidmat-relief-copilot  ·  relief-web-4idrhaffca-el.a.run.app",
             0.5, 7.05, 12.4, 0.35, size=11, color=MUTED, align=PP_ALIGN.LEFT)


def section_label(slide, text: str):
    add_text(slide, text.upper(), 0.7, 0.35, 5, 0.4, size=13, bold=True, color=CANOPY)


def big_title(slide, text: str, y: float = 1.0):
    add_text(slide, text, 0.7, y, 11.9, 1.6, size=40, bold=True, color=WHITE, font_name="Segoe UI Semibold")


def subtitle(slide, text: str, y: float = 2.7):
    add_text(slide, text, 0.7, y, 11.9, 0.9, size=20, color=MUTED)


# ── SLIDE BUILDERS ─────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)

    add_text(s, "ALKHIDMAT RELIEF COPILOT", 0.7, 1.2, 11.9, 0.6, size=16, bold=True, color=CANOPY)
    add_text(s, "From an aid request to a\nverified ticket — with a\nhuman when it matters.", 0.7, 1.9, 11.9, 2.8,
             size=48, bold=True, color=WHITE, font_name="Segoe UI Semibold")
    add_text(s, "Multi-agent NGO aid desk  ·  Urdu + English  ·  Qwen via DashScope  ·  Alibaba Cloud AI Hackathon Pakistan 2026",
             0.7, 4.9, 11.9, 0.6, size=16, color=MUTED)

    # Accent chips
    add_chip(s, 0.7, 5.8, 2.0, 0.5, "LangGraph", CANOPY, BG_DARK, 14)
    add_chip(s, 2.85, 5.8, 2.0, 0.5, "FastAPI", CANOPY, BG_DARK, 14)
    add_chip(s, 5.0, 5.8, 2.0, 0.5, "Next.js 14", CANOPY, BG_DARK, 14)
    add_chip(s, 7.15, 5.8, 2.4, 0.5, "Postgres/pgvector", CANOPY, BG_DARK, 14)
    add_chip(s, 9.7, 5.8, 2.6, 0.5, "Qwen / DashScope", BRASS, BG_DARK, 14)

    banner = os.path.join(os.path.dirname(__file__), "assets", "readme-banner.png")
    if os.path.exists(banner):
        s.shapes.add_picture(banner, Inches(8.5), Inches(1.0), width=Inches(4.5))

    footer(s)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "The Problem")
    big_title(s, "Relief requests arrive messy.\nCritical cases get buried.")
    add_multiline(s, [
        "Free-text Urdu/English — incomplete contact and location",
        "Wrong category routing (food vs medical vs shelter)",
        "Duplicate requests from the same phone or area",
        "High-risk medical cases need a human gate — fast",
        "The first 60 seconds after a request is the bottleneck",
    ], 0.7, 3.2, 11.9, 3.5, size=20, color=WHITE)
    footer(s)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Our Solution")
    big_title(s, "An agentic operations desk\nthat verifies, matches, and routes.")

    add_multiline(s, [
        "6-agent governed pipeline: Intake → Triage → Knowledge → Integrity → Matcher → Dispatch",
        "HITL supervisor gate for critical or high-risk cases",
        "Ticket ID + resource match + next-step visibility for everyone",
        "Urdu and English intake — Qwen via Alibaba DashScope",
        "SOP retrieval (pgvector RAG) with readable citations",
    ], 0.7, 3.2, 11.9, 3.5, size=20, color=WHITE)
    footer(s)


def slide_winning_signal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Winning Signal")
    big_title(s, "Everyone knows what happens next.")

    # Left column
    add_text(s, "After submission", 0.7, 3.1, 5.5, 0.4, size=16, bold=True, color=CANOPY)
    add_multiline(s, [
        "Ticket ID: AKD-YYYYMMDD-XXXXXX",
        "Resource match (inventory + volunteer)",
        "Or: waiting for supervisor approval",
        "Full agent trace + case timeline",
    ], 0.7, 3.6, 5.5, 3.0, size=18, color=WHITE)

    # Right column
    add_text(s, "Promise", 7.0, 3.1, 5.5, 0.4, size=16, bold=True, color=BRASS)
    add_multiline(s, [
        "One minute to verified ticket",
        "Integrity never skipped on create",
        "HITL decisions are traceable",
        "Citizens never need a staff login",
    ], 7.0, 3.6, 5.5, 3.0, size=18, color=WHITE, bullet_color=BRASS)
    footer(s)


def slide_pipeline_sequence(prs, active_idx: int):
    """One slide per pipeline stage — creates step-by-step animation feel."""
    stages = ["Intake", "Triage", "Knowledge", "Integrity", "HITL Gate", "Matcher", "Dispatch"]
    descs = [
        "Language detect · extract need, location, phone, name",
        "Classify: Food, Medical, Shelter, Blood, Education, Other + priority",
        "Retrieve Alkhidmat SOPs (pgvector RAG or keyword fallback)",
        "Duplicate phone check · risk score · fraud heuristics",
        "Pause graph — wait for supervisor approve / reject",
        "Match inventory + assign volunteer by category and area",
        "Create ticket AKD-… · notify requester · log audit",
    ]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, f"Pipeline  ·  Step {active_idx + 1} of {len(stages)}")
    big_title(s, stages[active_idx], y=0.9)
    add_text(s, descs[active_idx], 0.7, 2.0, 11.9, 0.5, size=20, color=MUTED)

    # Pipeline strip
    chip_w = 1.55
    gap = 0.18
    total_w = len(stages) * chip_w + (len(stages) - 1) * gap
    start_x = (13.333 - total_w) / 2
    y = 3.2

    for i, stage in enumerate(stages):
        is_active = i == active_idx
        is_done = i < active_idx
        if is_active:
            fill = CANOPY
            tc = BG_DARK
        elif is_done:
            fill = "#1a3d2e"
            tc = CANOPY
        else:
            fill = "#1a2420"
            tc = MUTED

        add_chip(s, start_x + i * (chip_w + gap), y, chip_w, 0.7, stage, fill, tc, 13, is_active)

        if i < len(stages) - 1:
            arrow_color = CANOPY if i < active_idx else MUTED
            add_text(s, "→", start_x + i * (chip_w + gap) + chip_w + 0.01, y + 0.12,
                     gap - 0.02, 0.5, size=22, bold=True, color=arrow_color, align=PP_ALIGN.CENTER)

    # Bottom detail card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(4.4), Inches(11.9), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(BG_CARD)
    card.line.color.rgb = rgb(DARK_LINE)
    card.line.width = Pt(1)

    detail_map = [
        ["Detects Urdu vs English automatically", "Extracts structured fields: need, location, phone, name",
         "Powered by Qwen via DashScope"],
        ["Classifies into 6 categories: Food, Medical, Shelter, Blood, Education, Other",
         "Assigns priority: low / medium / high / critical",
         "LLM-based classification with Qwen"],
        ["Retrieves relevant Alkhidmat SOPs for the category",
         "pgvector similarity search with keyword fallback",
         "Citations shown as readable purpose + bullet points"],
        ["Checks phone / area for duplicate requests",
         "Computes risk score with fraud heuristics",
         "Flags cases that need HITL supervisor approval"],
        ["Graph pauses via LangGraph interrupt_before",
         "Supervisor approves → pipeline resumes to Matcher",
         "Supervisor rejects → case closed with traceable decision"],
        ["Matches available inventory by category and area",
         "Assigns nearest volunteer from the pool",
         "Resource match stored on the case record"],
        ["Creates ticket: AKD-YYYYMMDD-XXXXXX",
         "Logs full audit trail + case events",
         "Requester sees ticket ID or 'waiting for supervisor'"],
    ]

    add_multiline(s, detail_map[active_idx], 1.1, 4.65, 11.2, 1.8, size=17, color=WHITE, bullet_color=CANOPY)
    footer(s)


def slide_hitl(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Human-In-The-Loop")
    big_title(s, "Critical cases pause for\nhuman approval.")

    add_multiline(s, [
        "Integrity flags duplicate / fraud / high-risk cases",
        "LangGraph interrupt_before pauses the graph at HITL gate",
        "Supervisor reviews case with full agent trace",
        "Approve → pipeline resumes to Matcher → Dispatch",
        "Reject → case closed with traceable decision and note",
        "Durable checkpoints survive restarts",
    ], 0.7, 3.1, 7.5, 3.5, size=18, color=WHITE)

    # Approve / Reject chips
    add_chip(s, 9.0, 3.3, 3.5, 0.85, "✓  Approve", CANOPY, BG_DARK, 20)
    add_chip(s, 9.0, 4.4, 3.5, 0.85, "✗  Reject", SAFFRON, WHITE, 20)
    add_text(s, "Governed by LangGraph\ncheckpoints + durable state", 9.0, 5.5, 3.5, 0.8,
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)


def slide_public_staff(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Product Design")
    big_title(s, "Citizens submit publicly.\nOperators use staff login.")

    # Two columns with cards
    # Left: public
    card_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.0), Inches(5.7), Inches(3.6))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = rgb(BG_CARD)
    card_l.line.color.rgb = rgb(CANOPY)
    card_l.line.width = Pt(1.5)

    add_text(s, "PUBLIC  ·  /request", 1.0, 3.15, 5.1, 0.4, size=14, bold=True, color=CANOPY)
    add_multiline(s, [
        "No account or password required",
        "Landing → Request aid → submit in Urdu or English",
        "Sees ticket ID or 'waiting for supervisor'",
        "Check status by ticket + phone on /status",
    ], 1.0, 3.65, 5.1, 2.6, size=16, color=WHITE)

    # Right: staff
    card_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(3.0), Inches(5.7), Inches(3.6))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = rgb(BG_CARD)
    card_r.line.color.rgb = rgb(BRASS)
    card_r.line.width = Pt(1.5)

    add_text(s, "STAFF  ·  /login (JWT)", 7.2, 3.15, 5.1, 0.4, size=14, bold=True, color=BRASS)
    add_multiline(s, [
        "Desk operator: tickets, dashboard, case detail, PDF",
        "Supervisor: approve/reject HITL + all desk views",
        "Agent trace + SOP citations + case timeline",
        "Role-based API gates (JWT HS256)",
    ], 7.2, 3.65, 5.1, 2.6, size=16, color=WHITE, bullet_color=BRASS)
    footer(s)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Architecture")
    big_title(s, "System design", y=0.85)

    # Layer boxes
    layers = [
        ("Frontend", "Next.js 14 App Router", CANOPY, 0.7, 2.5, 5.7, 0.9),
        ("Backend", "FastAPI + SSE streaming", CANOPY, 6.9, 2.5, 5.7, 0.9),
        ("Orchestration", "LangGraph multi-agent pipeline", CANOPY, 0.7, 3.7, 5.7, 0.9),
        ("Database", "PostgreSQL + pgvector", CANOPY, 6.9, 3.7, 5.7, 0.9),
        ("LLM + Embeddings", "Qwen via DashScope (Alibaba Cloud)", BRASS, 0.7, 4.9, 11.9, 0.9),
        ("Deploy", "GCP Cloud Run + Cloud SQL (live)", MUTED, 0.7, 6.05, 11.9, 0.65),
    ]

    for label, desc, color, x, y, w, h in layers:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(BG_CARD)
        card.line.color.rgb = rgb(color)
        card.line.width = Pt(1.5)
        add_text(s, label, x + 0.25, y + 0.1, w - 0.5, 0.3, size=13, bold=True, color=color)
        add_text(s, desc, x + 0.25, y + 0.42, w - 0.5, 0.35, size=16, color=WHITE)

    # Arrows
    add_text(s, "→", 6.35, 2.7, 0.5, 0.5, size=28, bold=True, color=CANOPY, align=PP_ALIGN.CENTER)
    add_text(s, "↓", 3.3, 3.45, 0.5, 0.3, size=22, bold=True, color=CANOPY, align=PP_ALIGN.CENTER)
    add_text(s, "↓", 9.6, 3.45, 0.5, 0.3, size=22, bold=True, color=CANOPY, align=PP_ALIGN.CENTER)

    footer(s)


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "Demo Walkthrough")
    big_title(s, "Citizen → Desk → HITL → Dispatch")

    add_multiline(s, [
        "Open landing (/) and click Request aid",
        "Submit an Urdu or English request — no account needed",
        "Watch live pipeline strip: each agent lights as it runs",
        "See ticket ID or 'Waiting for supervisor' above the fold",
        "Run a duplicate/critical case to trigger HITL",
        "Login as supervisor → approve → ticket dispatched",
        "Open case timeline → export PDF for ops record",
        "Check status as citizen: ticket + phone on /status",
    ], 0.7, 3.0, 11.9, 4.0, size=18, color=WHITE)
    footer(s)


def slide_outcomes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)
    section_label(s, "What We Deliver")
    big_title(s, "A governed aid desk,\nready for real operations.")

    # Left
    add_text(s, "Product capabilities", 0.7, 3.1, 6.0, 0.4, size=14, bold=True, color=CANOPY)
    add_multiline(s, [
        "Verified, prioritized ticket creation",
        "Duplicate detection + integrity risk score",
        "Resource matching + volunteer assignment",
        "Supervisor HITL for critical cases",
        "Agent trace + SOP citations + timeline",
        "PDF export for desk records",
    ], 0.7, 3.6, 6.0, 3.2, size=16, color=WHITE)

    # Right
    add_text(s, "Tech stack", 7.5, 3.1, 5.0, 0.4, size=14, bold=True, color=BRASS)
    add_multiline(s, [
        "LangGraph (orchestration + checkpoints)",
        "FastAPI + SSE real-time streaming",
        "PostgreSQL + pgvector (RAG)",
        "Qwen via DashScope (Alibaba Cloud)",
        "Next.js 14 (App Router)",
        "GCP Cloud Run + Cloud SQL (live)",
    ], 7.5, 3.6, 5.0, 3.2, size=16, color=WHITE, bullet_color=BRASS)
    footer(s)


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(s, BG_DARK)
    accent_bar(s)

    add_text(s, "ALKHIDMAT RELIEF COPILOT", 0.7, 2.0, 11.9, 0.5, size=16, bold=True, color=CANOPY)
    add_text(s, "From messy aid request to\nverified ticket — with a\nhuman when it matters.",
             0.7, 2.6, 11.9, 2.5, size=44, bold=True, color=WHITE, font_name="Segoe UI Semibold")

    add_text(s, "Muhammad Atif  ·  P00627  ·  Alibaba Cloud AI Hackathon Pakistan 2026",
             0.7, 5.3, 11.9, 0.5, size=16, color=MUTED)

    add_chip(s, 0.7, 6.0, 4.0, 0.6, "github.com/Atif1299/alkhidmat-relief-copilot", CANOPY, BG_DARK, 14, False)
    add_chip(s, 5.0, 6.0, 5.5, 0.6, "Live: relief-web-4idrhaffca-el.a.run.app", BRASS, BG_DARK, 14, False)

    footer(s)


# ── MAIN ────────────────────────────────────────────────────────────────

def generate():
    prs = Presentation()
    # Keep default widescreen 13.33 x 7.5

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_winning_signal(prs)

    # 7 pipeline sequence slides (the "animation")
    for i in range(7):
        slide_pipeline_sequence(prs, i)

    slide_hitl(prs)
    slide_public_staff(prs)
    slide_architecture(prs)
    slide_demo(prs)
    slide_outcomes(prs)
    slide_closing(prs)

    out = os.path.join(os.path.dirname(__file__), "submission_presentation.pptx")
    prs.save(out)
    print(f"Generated {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    generate()
